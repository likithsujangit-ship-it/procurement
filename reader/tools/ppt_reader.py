"""
PowerPoint (PPTX/PPT) Content Extraction Module.
Reads slides, text shapes, notes, and embedded images using python-pptx and pytesseract.
"""

import io
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tools.utils import setup_logger, preprocess_image_for_ocr

logger = setup_logger("ppt_reader")


def extract_ppt_text(filepath: Path) -> str:
    """
    Extracts text from slides in a .pptx file, and performs OCR on embedded pictures.
    Provides troubleshooting fallback for older .ppt files.
    
    Args:
        filepath: Path to the PowerPoint document on disk.
        
    Returns:
        Extracted text content as a string.
    """
    logger.info(f"Extracting text from PPTX: {filepath.name}")
    
    # Handle older .ppt files gracefully
    if filepath.suffix.lower() == ".ppt":
        logger.warning(f"File {filepath.name} is a legacy .ppt file. python-pptx only supports .pptx.")
        return (
            "[Unsupported Format: Legacy .ppt file]\n"
            "Troubleshooting: Please convert this file to the modern .pptx format "
            "using Microsoft PowerPoint or LibreOffice to allow automated content reading."
        )

    native_chars = 0
    ocr_chars = 0

    try:
        prs = Presentation(filepath)
        text_parts = []
        
        try:
            import pytesseract
            from PIL import Image
            has_ocr = True
        except ImportError:
            has_ocr = False
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = [f"--- Slide {slide_idx + 1} ---"]
            slide_native_parts = []
            
            # Extract text and embedded images from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_native_parts.append(shape.text.strip())
                    native_chars += len(shape.text.strip())
                    
                if has_ocr:
                    # Check for picture shapes or objects with image attribute
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                        try:
                            image_bytes = shape.image.blob
                            pil_img = Image.open(io.BytesIO(image_bytes))
                            processed_img = preprocess_image_for_ocr(pil_img)
                            ocr_text = pytesseract.image_to_string(processed_img).strip()
                            if ocr_text:
                                ocr_chars += len(ocr_text)
                                slide_content.append(f"[OCR from embedded image, slide {slide_idx + 1}]: {ocr_text}")
                        except Exception as img_err:
                            logger.debug(f"Failed to process image on slide {slide_idx + 1}: {img_err}")
                            
            if slide_native_parts:
                slide_content.insert(1, "\n".join(slide_native_parts))
            elif len(slide_content) == 1:
                slide_content.append("[No text on slide]")
                
            text_parts.append("\n".join(slide_content))
                
        logger.debug(f"PPTX Extraction Audit for {filepath.name}: native_chars={native_chars}, ocr_chars={ocr_chars}")

        if not text_parts:
            return "[PowerPoint presentation has no slides]"
            
        return "\n\n".join(text_parts)

    except Exception as e:
        logger.error(f"Error reading PPTX {filepath.name}: {e}")
        return f"[Error reading PPTX file: {e}]"
